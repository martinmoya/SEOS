"""
Document Processing Service.
Orchestrates document AI operations preserving format.
"""

from pathlib import Path


class DocumentProcessingService:
    def __init__(self, processor):
        self.processor = processor

    def process(self, source: Path, output_suffix: str) -> Path:
        suffix = source.suffix.lower()

        # Forzar que los documentos procesados se guarden en projects/
        dest_dir = Path("projects")
        dest_dir.mkdir(parents=True, exist_ok=True)
        destination = dest_dir / f"{source.stem}.{output_suffix}{source.suffix}"

        if suffix in (".txt", ".md"):
            self._process_text(source, destination)
        elif suffix == ".pdf":
            self._process_pdf(source, destination)
        elif suffix == ".docx":
            self._process_docx(source, destination)
        elif suffix == ".xlsx":
            self._process_xlsx(source, destination)
        elif suffix == ".pptx":
            self._process_pptx(source, destination)
        else:
            raise ValueError(f"Unsupported format: {suffix}")

        return destination

    def _process_text(self, source: Path, dest: Path):
        text = source.read_text(encoding="utf-8")
        translated = self.processor.process(text)
        dest.write_text(translated, encoding="utf-8")

    def _process_pdf(self, source: Path, dest: Path):
        # Usamos el nuevo lector avanzado con OCR
        from documents.pdf_reader import PdfReaderDocument
        from reportlab.platypus import SimpleDocTemplate, Paragraph
        from reportlab.lib.styles import getSampleStyleSheet

        reader = PdfReaderDocument()
        text = reader.read(source)

        # Si el lector no pudo extraer nada (ni con OCR), abortamos
        if not text.strip():
            raise ValueError(
                "Could not extract text from PDF. OCR might have failed or file is empty."
            )

        translated = self.processor.process(text)

        styles = getSampleStyleSheet()
        story = []
        for line in translated.splitlines():
            if line.strip():
                story.append(Paragraph(line, styles["Normal"]))

        pdf = SimpleDocTemplate(str(dest))
        pdf.build(story)

    def _process_docx(self, source: Path, dest: Path):
        from docx import Document

        doc = Document(str(source))

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                translated = self.processor.process(paragraph.text)
                if paragraph.runs:
                    paragraph.runs[0].text = translated
                    for run in paragraph.runs[1:]:
                        run.text = ""
                else:
                    paragraph.text = translated

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():
                            translated = self.processor.process(paragraph.text)
                            if paragraph.runs:
                                paragraph.runs[0].text = translated
                                for run in paragraph.runs[1:]:
                                    run.text = ""
                            else:
                                paragraph.text = translated

        doc.save(str(dest))

    def _process_xlsx(self, source: Path, dest: Path):
        from openpyxl import load_workbook

        wb = load_workbook(str(source))

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if (
                        cell.value
                        and isinstance(cell.value, str)
                        and cell.value.strip()
                    ):
                        cell.value = self.processor.process(cell.value)

        wb.save(str(dest))

    def _process_pptx(self, source: Path, dest: Path):
        from pptx import Presentation

        prs = Presentation(str(source))

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                run.text = self.processor.process(run.text)

        prs.save(str(dest))
