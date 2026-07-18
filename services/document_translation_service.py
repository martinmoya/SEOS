"""
Document Translation Service.
Orchestrates document translation preserving format where possible.
"""

from pathlib import Path
from processors.translation_processor import TranslationProcessor


class DocumentTranslationService:
    def __init__(self, processor: TranslationProcessor):
        self.processor = processor

    def translate(self, source: Path, language: str) -> Path:
        suffix = source.suffix.lower()
        destination = source.with_suffix(f".{language}{source.suffix}")

        if suffix in (".txt", ".md"):
            self._translate_text(source, destination)
        elif suffix == ".pdf":
            self._translate_pdf(source, destination)
        elif suffix == ".docx":
            self._translate_docx(source, destination)
        elif suffix == ".xlsx":
            self._translate_xlsx(source, destination)
        elif suffix == ".pptx":
            self._translate_pptx(source, destination)
        else:
            raise ValueError(f"Unsupported format: {suffix}")

        return destination

    def _translate_text(self, source: Path, dest: Path):
        text = source.read_text(encoding="utf-8")
        translated = self.processor.process(text)
        dest.write_text(translated, encoding="utf-8")

    def _translate_pdf(self, source: Path, dest: Path):
        from pypdf import PdfReader
        from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet

        reader = PdfReader(source)
        styles = getSampleStyleSheet()
        story = []

        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                translated = self.processor.process(text)
                for line in translated.splitlines():
                    if line.strip():
                        story.append(Paragraph(line, styles["Normal"]))
            story.append(PageBreak())

        if story and isinstance(story[-1], PageBreak):
            story.pop()

        pdf = SimpleDocTemplate(str(dest))
        pdf.build(story)

    def _translate_docx(self, source: Path, dest: Path):
        from docx import Document

        doc = Document(str(source))

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                translated = self.processor.process(paragraph.text)
                # Preserve formatting by replacing text in the first run
                # and clearing the rest.
                if paragraph.runs:
                    paragraph.runs[0].text = translated
                    for run in paragraph.runs[1:]:
                        run.text = ""
                else:
                    paragraph.text = translated

        # Translate tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():
                            translated = self.processor.process(paragraph.text)
                            if paragraph.runs:
                                paragraph.runs[0].text = translated
                                for run in paragraph.runs[1:]:
                                    run.text = ""
                            else:
                                paragraph.text = translated

        doc.save(str(dest))

    def _translate_xlsx(self, source: Path, dest: Path):
        from openpyxl import load_workbook

        wb = load_workbook(str(source))

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if (
                        cell.value
                        and isinstance(cell.value, str)
                        and cell.value.strip()
                    ):
                        cell.value = self.processor.process(cell.value)

        wb.save(str(dest))

    def _translate_pptx(self, source: Path, dest: Path):
        from pptx import Presentation

        prs = Presentation(str(source))

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # Translate run by run to preserve character formatting
                        for run in paragraph.runs:
                            if run.text.strip():
                                run.text = self.processor.process(run.text)

        prs.save(str(dest))
